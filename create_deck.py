"""
Generate the Agentic SDLC at Schwab – Team Overview slide deck.
Uses the SchwabPlan.com reference PPTX as a template for the Mphasis layout/theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE = r"C:\Users\banu.parasuraman\Downloads\mainframe-modernization\SchwabPlan.com API Modernization-V3.pptx"
OUTPUT   = r"C:\Users\banu.parasuraman\Downloads\mainframe-modernization\Agentic_SDLC_at_Schwab_Overview.pptx"

# ── Brand colours (Mphasis / Schwab palette) ──────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5C)   # Schwab navy
MID_BLUE   = RGBColor(0x00, 0x6B, 0xA6)   # accent
LIGHT_BLUE = RGBColor(0x00, 0x96, 0xD6)   # highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x00, 0x7A, 0x33)
ORANGE     = RGBColor(0xE8, 0x6C, 0x00)
TEAL       = RGBColor(0x00, 0x80, 0x80)

# ── Helpers ────────────────────────────────────────────────────────────────────
def _font(run, name="Segoe UI", size=18, bold=False, color=DARK_GRAY):
    run.font.name  = name
    run.font.size  = Pt(size)
    run.font.bold  = bold
    if color:
        run.font.color.rgb = color

def _add_run(para, text, name="Segoe UI", size=18, bold=False, color=DARK_GRAY):
    run = para.add_run()
    run.text = text
    _font(run, name, size, bold, color)
    return run

def _set_cell(cell, text, size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    _add_run(p, text, size=size, bold=bold, color=color)
    cell.text_frame.word_wrap = True

def _header_cell(cell, text, size=13):
    _set_cell(cell, text, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

def _alt_row(cell, row_idx):
    if row_idx % 2 == 0:
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY

def _add_title_bar(slide, text, top=Inches(0.0), left=Inches(0.0),
                   width=Inches(13.33), height=Inches(0.95)):
    """Dark-blue title bar spanning the top."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = DARK_BLUE
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, "  " + text, size=28, bold=True, color=WHITE)
    return shp

def _add_subtitle(slide, text, top=Inches(1.1)):
    txBox = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.5))
    tf = txBox.text_frame
    _add_run(tf.paragraphs[0], text, size=20, bold=True, color=MID_BLUE)
    return txBox

def _add_body_text(slide, text, top=Inches(1.7), left=Inches(0.6),
                   width=Inches(12), height=Inches(5.2), size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        _add_run(p, line, size=size, color=DARK_GRAY)
        p.space_after = Pt(6)
    return txBox

def _add_bullet_list(slide, items, top=Inches(1.7), left=Inches(0.6),
                     width=Inches(11.5), height=Inches(5.2), size=16,
                     title_items=None):
    """Add a bulleted list. items = list of strings (or (bold_prefix, rest))."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.level = 0
        # bullet char
        if isinstance(item, tuple):
            _add_run(p, "●  " + item[0], size=size, bold=True, color=DARK_BLUE)
            _add_run(p, " – " + item[1], size=size, color=DARK_GRAY)
        else:
            _add_run(p, "●  " + item, size=size, color=DARK_GRAY)
    return txBox

def _add_table(slide, headers, rows, left=Inches(0.5), top=Inches(1.7),
               width=Inches(12.3), row_height=Inches(0.45)):
    cols = len(headers)
    num_rows = len(rows) + 1
    tbl_height = row_height * num_rows
    table_shape = slide.shapes.add_table(num_rows, cols, left, top, width, tbl_height)
    table = table_shape.table

    # distribute column widths
    col_w = int(width / cols)
    for ci in range(cols):
        table.columns[ci].width = col_w

    # header
    for ci, h in enumerate(headers):
        _header_cell(table.cell(0, ci), h)

    # data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            bold = (ci == 0)
            _set_cell(cell, val, size=13, bold=bold, color=DARK_GRAY)
            _alt_row(cell, ri)

    return table_shape

def _add_colored_box(slide, text, left, top, width, height, fill_color, text_color=WHITE, size=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _add_run(tf.paragraphs[0], text, size=size, bold=True, color=text_color)
    shp.text_frame.margin_left = Pt(6)
    shp.text_frame.margin_right = Pt(6)
    return shp

def _add_arrow(slide, left, top, width, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = MID_BLUE
    shp.line.fill.background()
    return shp

def _add_footer(slide, date_str="March 2026"):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.35))
    tf = txBox.text_frame
    _add_run(tf.paragraphs[0], f"© Mphasis 2026  |  Proprietary and confidential  |  {date_str}",
             size=10, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD THE DECK
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation(TEMPLATE)

# Remove all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

# Identify blank-ish layout (6_Custom Layout is index 2)
blank_layout = prs.slide_layouts[2]  # 6_Custom Layout
cover_layout = prs.slide_layouts[0]  # Cover only Mphasis

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – TITLE / COVER
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(cover_layout)
# Clear existing placeholder text
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

# Title
for ph in slide.placeholders:
    if "Title" in ph.name:
        p = ph.text_frame.paragraphs[0]
        _add_run(p, "Agentic SDLC at Schwab", size=36, bold=True, color=WHITE)
    elif "Subtitle" in ph.name:
        p = ph.text_frame.paragraphs[0]
        _add_run(p, "Legacy Modernization – Team Overview", size=22, bold=False, color=WHITE)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – AGENDA
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
# remove inherited placeholders text
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Agenda")
items = [
    "Executive Summary – Agentic SDLC Vision",
    "Portfolio of Modernization Programs",
    "Key Enablers: Vertex AI Search & Spec-Driven Development",
    "Program 1 – Bank Sweep: COBOL → Spring (Legacy Mainframe Modernization)",
    "Program 2 – RPS / SchwabPlan.com: .NET → .NET Core (LangChain/LangGraph & Google ADK)",
    "Program 3 – O2: Legacy Batch Modernization – COBOL → Spring Batch (Custom Agent)",
    "Agentic SDLC Architecture & Agent Framework Comparison",
    "Metrics & Efficiency Gains",
    "Why Mphasis",
]
_add_bullet_list(slide, items, top=Inches(1.3), size=18)
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – EXECUTIVE SUMMARY
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Executive Summary – Agentic SDLC Vision")

body = (
    "Schwab is modernizing multiple legacy estates—mainframe COBOL, .NET Framework, and legacy batch systems—"
    "using an AI-driven, agentic software development lifecycle (Agentic SDLC).\n\n"
    "Rather than a traditional manual rewrite, Mphasis employs autonomous AI agents that perform reverse engineering, "
    "requirements extraction, architecture design, code generation, testing, and deployment with human-in-the-loop "
    "governance at critical gates.\n\n"
    "Three flagship programs demonstrate this capability across diverse technology stacks:\n"
    "    1) Bank Sweep – COBOL to Spring (Mainframe Modernization)\n"
    "    2) RPS / SchwabPlan.com – .NET to .NET Core (LangChain/LangGraph & Google ADK)\n"
    "    3) O2 – Legacy Batch COBOL to Spring Batch (Custom Agent)\n\n"
    "Two cross-cutting enablers amplify every program:\n"
    "    ● Vertex AI Search – Domain knowledge grounding via enterprise search over legacy artifacts\n"
    "    ● Spec-Driven Development – Agents consume and produce formal specifications (OpenAPI, AsyncAPI, BRDs, "
    "Gherkin) as the single source of truth, ensuring accuracy and traceability."
)
_add_body_text(slide, body, top=Inches(1.2), size=16)
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – PORTFOLIO OVERVIEW
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Portfolio of Agentic SDLC Modernization Programs")

headers = ["Program", "Source Stack", "Target Stack", "Agent Framework", "Status"]
rows = [
    ["Bank Sweep",         "COBOL / Mainframe",         "Spring Boot (Java)",    "Custom Agentic Pipeline",       "In Progress"],
    ["RPS – SchwabPlan.com\n(LangChain/LangGraph)", ".NET Framework 4.8 (C#)", ".NET Core (C#)",  "LangChain / LangGraph",         "POC Complete"],
    ["RPS – SchwabPlan.com\n(Google ADK)",           ".NET Framework 4.8 (C#)", ".NET Core (C#)",  "Google ADK (Vertex AI)",        "POC Complete"],
    ["O2 Batch",           "COBOL Batch Jobs",          "Spring Batch (Java)",   "Custom Agent",                  "In Progress"],
]
_add_table(slide, headers, rows, top=Inches(1.3), row_height=Inches(0.65))

# Key callout
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(12), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Cross-Cutting Enablers", size=18, bold=True, color=DARK_BLUE)
p2 = tf.add_paragraph()
_add_run(p2, "●  Vertex AI Search", size=15, bold=True, color=MID_BLUE)
_add_run(p2, " – Enterprise search over legacy code, documentation, and domain knowledge stored in GCS. "
         "Provides grounded context to every agent so outputs reflect real business rules.", size=15, color=DARK_GRAY)
p3 = tf.add_paragraph()
_add_run(p3, "●  Spec-Driven Development", size=15, bold=True, color=MID_BLUE)
_add_run(p3, " – Agents consume and produce formal specs (OpenAPI, Gherkin, BRDs, DDD artifacts). "
         "Specs serve as contracts between agent phases, enabling validation, traceability, and repeatability.", size=15, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – KEY ENABLER: VERTEX AI SEARCH
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Key Enabler: Vertex AI Search for Domain Knowledge")

# Left description
body_left = (
    "Vertex AI Search provides an enterprise-grade, managed search "
    "and retrieval layer over Schwab's legacy artifacts:\n\n"
    "●  Legacy source code (COBOL, .NET, batch JCL)\n"
    "●  Business requirements documents & runbooks\n"
    "●  Database schemas, copybooks, data dictionaries\n"
    "●  Existing test cases and deployment manifests\n"
    "●  Domain glossaries and business-rule catalogs\n\n"
    "All artifacts are indexed in Google Cloud Storage (GCS) and made "
    "searchable via Vertex AI Search datastores."
)
_add_body_text(slide, body_left, top=Inches(1.2), left=Inches(0.5),
               width=Inches(5.8), height=Inches(4.0), size=15)

# Right – How agents use it
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "How Agents Leverage Vertex AI Search", size=17, bold=True, color=DARK_BLUE)

benefits = [
    ("Grounded Code Analysis", "Agents query Vertex AI Search to retrieve relevant legacy source, copybooks, and data models before analyzing code."),
    ("Contextual BRD Generation", "Agents pull business-rule documentation and runbooks to produce accurate BRDs without hallucination."),
    ("Domain-Aware Design", "DDD artifacts and architecture decisions are informed by real domain glossaries and bounded-context definitions."),
    ("Spec Validation", "Generated OpenAPI specs and Gherkin scenarios are cross-checked against indexed legacy test cases and requirements."),
    ("Continuous Learning", "As agents produce new artifacts, those artifacts are re-ingested into the search index, enriching future queries."),
]
for label, desc in benefits:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    _add_run(p2, f"●  {label}", size=14, bold=True, color=MID_BLUE)
    p3 = tf.add_paragraph()
    _add_run(p3, f"    {desc}", size=13, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – KEY ENABLER: SPEC-DRIVEN DEVELOPMENT
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Key Enabler: Spec-Driven Development")

body_left = (
    "In Spec-Driven Development, formal specifications are the single source "
    "of truth at every stage of the Agentic SDLC.\n\n"
    "Instead of free-form prose, agents produce and consume structured artifacts:\n\n"
    "  ● OpenAPI / AsyncAPI specs for API contracts\n"
    "  ● Gherkin (BDD) scenarios for acceptance criteria\n"
    "  ● Domain-Driven Design models (aggregates, bounded contexts)\n"
    "  ● Business Requirements Documents (BRDs) with traceable rules\n"
    "  ● Architecture Decision Records (ADRs)\n\n"
    "Each specification becomes a contract between agent phases. "
    "The next agent validates its input spec before proceeding, "
    "creating a self-correcting, auditable pipeline."
)
_add_body_text(slide, body_left, top=Inches(1.2), left=Inches(0.5),
               width=Inches(5.8), height=Inches(5.0), size=15)

# Right – Flow diagram (text representation)
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Spec-Driven Agent Pipeline", size=17, bold=True, color=DARK_BLUE)

steps = [
    ("1. Legacy Intake", "Agent ingests source code → produces Code Analysis Spec"),
    ("2. Requirements", "Agent consumes Code Analysis Spec → produces BRD + Gherkin"),
    ("3. Domain Design", "Agent consumes BRD → produces DDD Artifacts + Architecture Spec"),
    ("4. API Design", "Agent consumes DDD Artifacts → produces OpenAPI / AsyncAPI Spec"),
    ("5. Code Generation", "Agent consumes OpenAPI Spec → produces Target Code + Unit Tests"),
    ("6. Validation", "Agent consumes Gherkin + Target Code → runs tests, feeds back"),
    ("7. Deployment", "Agent consumes validated artifacts → CI/CD pipeline execution"),
]
for step, desc in steps:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    _add_run(p2, step, size=14, bold=True, color=MID_BLUE)
    p3 = tf.add_paragraph()
    _add_run(p3, f"    {desc}", size=13, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – PROGRAM 1: BANK SWEEP (COBOL → Spring)
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Program 1 – Bank Sweep: COBOL → Spring Boot")
_add_subtitle(slide, "Legacy Mainframe Modernization", top=Inches(1.1))

items = [
    ("Source", "COBOL programs running on mainframe with JCL batch orchestration"),
    ("Target", "Spring Boot microservices (Java 17+) deployed on PCF / Kubernetes"),
    ("Agent Framework", "Custom Agentic pipeline leveraging Vertex AI (Gemini models)"),
    ("Vertex AI Search", "Indexes COBOL copybooks, JCL, data dictionaries, and business-rule catalogs from GCS to ground agent analysis"),
    ("Spec-Driven Flow", "Agents produce BRDs from COBOL analysis → Gherkin acceptance criteria → DDD bounded contexts → OpenAPI specs → Spring Boot code"),
    ("Human-in-the-Loop", "SME review gates after BRD extraction and before code generation"),
    ("Key Outcome", "Automated reverse engineering of COBOL business rules with high accuracy; repeatable factory model for additional mainframe modules"),
]
_add_bullet_list(slide, items, top=Inches(1.7), size=15)
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – PROGRAM 2: RPS / SCHWABPLAN.COM
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Program 2 – RPS / SchwabPlan.com: .NET → .NET Core")
_add_subtitle(slide, "Legacy .NET Modernization with Dual Agent Frameworks", top=Inches(1.1))

# Two-column approach boxes
_add_colored_box(slide, "LangChain / LangGraph Approach", Inches(0.5), Inches(1.75),
                 Inches(5.8), Inches(0.55), MID_BLUE, WHITE, 16)
lang_items = (
    "●  Multi-model LLM runtime (Gemini + fallback)\n"
    "●  Graph-based orchestration via LangGraph\n"
    "●  Shared memory & state across agent nodes\n"
    "●  Tool / MCP integrations (GitHub, Jira, Confluence)\n"
    "●  Governance & observability via LangSmith\n"
    "●  Complex feedback loops & conditional branching\n"
    "●  Spec validation at each graph node transition"
)
_add_body_text(slide, lang_items, top=Inches(2.4), left=Inches(0.6),
               width=Inches(5.6), height=Inches(3.0), size=14)

_add_colored_box(slide, "Google ADK Approach", Inches(6.8), Inches(1.75),
                 Inches(5.8), Inches(0.55), TEAL, WHITE, 16)
adk_items = (
    "●  GCP-native, Gemini-centric enterprise stack\n"
    "●  Agent abstractions with roles + tools\n"
    "●  Multi-agent orchestration (Vertex AI)\n"
    "●  Shared memory & state management\n"
    "●  Tool / MCP integrations (GitHub, Jira, CI/CD)\n"
    "●  Governance & observability built-in\n"
    "●  Spec-driven validation agents per stage"
)
_add_body_text(slide, adk_items, top=Inches(2.4), left=Inches(6.9),
               width=Inches(5.6), height=Inches(3.0), size=14)

# Bottom callout
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Both approaches leverage:", size=15, bold=True, color=DARK_BLUE)
p2 = tf.add_paragraph()
_add_run(p2, "●  Vertex AI Search for domain knowledge grounding  ", size=14, bold=False, color=DARK_GRAY)
_add_run(p2, "●  Spec-Driven Development for artifact traceability  ", size=14, bold=False, color=DARK_GRAY)
_add_run(p2, "●  Human-in-the-loop at validation gates", size=14, bold=False, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – RPS CODEBASE METRICS (from reference deck)
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "RPS Codebase Metrics – .NET Framework 4.8 (C#)")

headers = ["Metric", "Schwabplan", "Schwabplan2.net", "EDP", "Loans Web App", "EDP Loan"]
rows = [
    ["Files",           "6,067",  "3,058",  "1,881",  "12",   "70"],
    ["Classes",         "10,616", "5,297",  "3,551",  "12",   "146"],
    ["Methods",         "62,262", "30,946", "13,213", "45",   "766"],
    ["Business Rules",  "24,990", "12,510", "6,283",  "61",   "302"],
    ["Function Rules",  "1,235",  "620",    "77",     "3",    "2"],
]
_add_table(slide, headers, rows, top=Inches(1.3), row_height=Inches(0.55))

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Overall Scale", size=16, bold=True, color=DARK_BLUE)
p2 = tf.add_paragraph()
_add_run(p2, "~7,948 files  |  ~14,167 classes  |  ~75K methods  |  ~33K business rules  |  ~1,312 function-level rules",
         size=15, color=DARK_GRAY)
p3 = tf.add_paragraph()
_add_run(p3, "→ Large, rules-heavy monolithic estate requiring systematic, agent-driven modernization",
         size=15, bold=False, color=MID_BLUE)
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – PROGRAM 3: O2 BATCH MODERNIZATION
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Program 3 – O2: Legacy Batch Modernization")
_add_subtitle(slide, "COBOL Batch → Spring Batch (Custom Agent)", top=Inches(1.1))

items = [
    ("Source", "COBOL batch programs with JCL scheduling, sequential file I/O, and DB2 interactions"),
    ("Target", "Spring Batch (Java) with modern job scheduling, chunk-oriented processing, and cloud-ready deployment"),
    ("Agent Framework", "Custom-built agent specifically designed for batch-to-batch modernization patterns"),
    ("Vertex AI Search", "Indexes JCL procedures, COBOL copybooks, file layouts, and batch run documentation to provide agents with scheduling dependencies and data-flow context"),
    ("Spec-Driven Flow", "Agent produces Batch Job Specs from JCL/COBOL analysis → Step-level BRDs → Spring Batch Job configuration → Chunk reader/processor/writer code → Integration tests"),
    ("Key Differentiator", "Custom agent handles batch-specific patterns: checkpoint/restart logic, commit intervals, conditional step execution, and file-based partitioning"),
    ("Human-in-the-Loop", "Mandatory review of batch job dependency graphs and data-flow mappings before code generation"),
]
_add_bullet_list(slide, items, top=Inches(1.7), size=15)
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – AGENTIC SDLC ARCHITECTURE
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Agentic SDLC Architecture – Using Schwab Ecosystem")

# Pipeline stages as colored boxes
stages = [
    ("Legacy\nIntake", DARK_BLUE),
    ("Code\nAnalysis", MID_BLUE),
    ("Requirements\n(BRD)", TEAL),
    ("DDD\nArtifacts", GREEN),
    ("API Design\n(OpenAPI)", MID_BLUE),
    ("User\nStories", DARK_BLUE),
    ("Code\nGeneration", TEAL),
    ("Testing &\nValidation", GREEN),
    ("Deployment\n(CI/CD)", DARK_BLUE),
]

box_w = Inches(1.25)
box_h = Inches(0.85)
start_left = Inches(0.3)
top_row = Inches(1.3)

for i, (label, color) in enumerate(stages):
    left = start_left + i * Inches(1.42)
    _add_colored_box(slide, label, left, top_row, box_w, box_h, color, WHITE, 11)
    if i < len(stages) - 1:
        _add_arrow(slide, left + box_w + Inches(0.02), top_row + Inches(0.28), Inches(0.14), Inches(0.28))

# Vertex AI Search bar
_add_colored_box(slide, "Vertex AI Search – Domain Knowledge (GCS)", Inches(0.3), Inches(2.35),
                 Inches(12.7), Inches(0.5), ORANGE, WHITE, 14)

# Spec-Driven bar
_add_colored_box(slide, "Spec-Driven Development – Formal Specs as Contracts Between Stages", Inches(0.3), Inches(2.95),
                 Inches(12.7), Inches(0.5), RGBColor(0x6B, 0x21, 0xA8), WHITE, 14)

# Three framework columns
col_w = Inches(3.9)
col_top = Inches(3.7)

# LangChain/LangGraph
_add_colored_box(slide, "LangChain / LangGraph", Inches(0.3), col_top, col_w, Inches(0.45), MID_BLUE, WHITE, 14)
lang_details = (
    "● Graph-based orchestration\n"
    "● Multi-model LLM runtime\n"
    "● Conditional branching & loops\n"
    "● LangSmith observability\n"
    "● Used for: RPS / SchwabPlan.com"
)
_add_body_text(slide, lang_details, top=col_top + Inches(0.55), left=Inches(0.4),
               width=Inches(3.7), height=Inches(2.5), size=13)

# Google ADK
_add_colored_box(slide, "Google ADK (Vertex AI)", Inches(4.5), col_top, col_w, Inches(0.45), TEAL, WHITE, 14)
adk_details = (
    "● GCP-native, Gemini-centric\n"
    "● Role-based agent abstractions\n"
    "● Multi-agent orchestration\n"
    "● Built-in governance\n"
    "● Used for: RPS / SchwabPlan.com"
)
_add_body_text(slide, adk_details, top=col_top + Inches(0.55), left=Inches(4.6),
               width=Inches(3.7), height=Inches(2.5), size=13)

# Custom Agent
_add_colored_box(slide, "Custom Agent", Inches(8.7), col_top, col_w, Inches(0.45), DARK_BLUE, WHITE, 14)
custom_details = (
    "● Purpose-built pipelines\n"
    "● Batch-specific patterns\n"
    "● COBOL/JCL specialization\n"
    "● Checkpoint/restart logic\n"
    "● Used for: Bank Sweep, O2"
)
_add_body_text(slide, custom_details, top=col_top + Inches(0.55), left=Inches(8.8),
               width=Inches(3.7), height=Inches(2.5), size=13)

# Tool/MCP integration bar
_add_colored_box(slide, "MCP Integrations: GitHub  |  Jira / Confluence  |  CI/CD  |  PCF  |  Selenium  |  Observability",
                 Inches(0.3), Inches(6.4), Inches(12.7), Inches(0.45), MED_GRAY, WHITE, 12)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – AGENT FRAMEWORK COMPARISON
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Agent Framework Comparison")

headers = ["Dimension", "VS Code + Copilot\n(Human-Assisted)", "Google ADK\n(Agent-Orchestrated)", "LangGraph\n(Graph-Orchestrated)", "Custom Agent\n(Batch-Specialized)"]
rows = [
    ["Primary Role",       "Developer productivity\naccelerator",     "Enterprise-grade agentic\nautomation",  "Advanced multi-agent\nworkflow orchestration", "Batch modernization\nspecialist"],
    ["Human Involvement",  "High (every step)",                       "Medium (at gates)",                     "Medium–Low (checkpoints)",                     "Medium (at gates)"],
    ["Automation Level",   "Low–Medium",                              "High",                                  "High",                                         "High"],
    ["Repeatability",      "Developer-dependent",                     "High (policy-driven)",                  "High (graph/state-driven)",                    "High (pattern-driven)"],
    ["Validation Model",   "Manual review +\nCopilot suggestions",    "Embedded validation\nagents per stage", "Embedded validation\nnodes per graph step",    "Custom validation\nper batch step"],
    ["Vertex AI Search",   "Partial (manual query)",                  "Integrated via\nSearch MCP",            "Integrated via\nSearch MCP",                   "Integrated via\nsearch tooling"],
    ["Spec-Driven",        "Partial",                                 "Full (specs as contracts)",             "Full (specs as contracts)",                    "Full (batch job specs)"],
    ["Best Fit",           "Small teams,\nearly pilots",              "Large portfolios,\nregulated envs",     "Complex flows,\ncross-system",                 "Batch / mainframe\nmodernization"],
    ["Schwab Programs",    "All (developer aid)",                     "RPS / SchwabPlan.com",                  "RPS / SchwabPlan.com",                         "Bank Sweep, O2"],
]
_add_table(slide, headers, rows, top=Inches(1.15), row_height=Inches(0.58))
_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 – LEVERAGING SCHWAB AI TOOLS / METHODOLOGY
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Leveraging Schwab's AI Tools & Methodology")

headers2 = ["Modernization\nMethod", "Legacy\nIntake", "Code\nAnalysis", "Requirements", "DDD\nArtifacts",
            "Design Doc", "User Stories", "API Spec", "Code\nGeneration", "Deployment"]
rows2 = [
    ["VS Code + Copilot",  "Git/GitHub\nMCP", "Human +\nAgent", "Human +\nAgent", "Human +\nAgent",
     "Human +\nAgent", "Human +\nAgent", "Human +\nAgent", "Human +\nAgent", "CI/CD\nAutomation"],
    ["Google ADK Agent",   "Git/GitHub\nMCP", "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation",
     "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation", "CI/CD\nAutomation"],
    ["LangGraph Agent",    "Git/GitHub\nMCP", "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation",
     "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation", "CI/CD\nAutomation"],
    ["Custom Agent\n(Batch)", "Git/GitHub\nMCP", "Agent +\nValidation", "Agent +\nValidation", "Agent +\nValidation",
     "Agent +\nValidation", "Agent +\nValidation", "Batch Job\nSpec", "Agent +\nValidation", "CI/CD\nAutomation"],
]
_add_table(slide, headers2, rows2, top=Inches(1.15), row_height=Inches(0.7))

# Time callout
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Time Comparison (per module):", size=16, bold=True, color=DARK_BLUE)

p2 = tf.add_paragraph()
_add_run(p2, "Manual Effort: ~200 hrs   |   ", size=15, color=ORANGE, bold=True)
_add_run(p2, "Human-in-Loop (Copilot): ~30 min   |   ", size=15, color=MID_BLUE, bold=True)
_add_run(p2, "End-to-End Agentic: ~15 min", size=15, color=GREEN, bold=True)

p3 = tf.add_paragraph()
p3.space_before = Pt(6)
_add_run(p3, "All methods leverage Vertex AI Search for domain grounding and Spec-Driven Development for artifact traceability.",
         size=14, color=MED_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 – METRICS & EFFICIENCY GAINS
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Metrics & Efficiency Gains (RPS POC)")

headers = ["Task", "Manual Effort (hrs)", "AI Effort (Prompts + Validation)", "Efficiency Gain %", "Accuracy"]
rows = [
    ["Relearn Code",                "50",  "10", "80%",    "High"],
    ["Generate BRD",                "32",  "6",  "81%",    "High"],
    ["Detailed BRDs",               "80",  "16", "80%",    "High"],
    ["Gherkin Scenarios",           "60",  "8",  "87%",    "High"],
    ["DDD Artifacts",               "40",  "7",  "83%",    "High"],
    ["API Design / Clean Arch",     "60",  "12", "80%",    "High"],
    ["User Story Generation",       "70",  "12", "83%",    "High"],
    ["Target Code Generation",      "200", "48", "76%",    "Medium"],
]
_add_table(slide, headers, rows, top=Inches(1.15), row_height=Inches(0.55))

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Key Takeaway: ", size=15, bold=True, color=DARK_BLUE)
_add_run(p, "76–87% efficiency gains across all SDLC phases. Vertex AI Search grounding and "
         "Spec-Driven Development contribute to High accuracy in analysis/design stages.",
         size=15, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 15 – WHY MPHASIS
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

_add_title_bar(slide, "Why Mphasis")

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Mphasis combines deep domain knowledge with a proven, governed AI modernization capability—"
         "allowing Schwab to modernize faster, safer, and at scale.", size=17, bold=False, color=MID_BLUE)

reasons = [
    ("1. Longstanding Partnership with Schwab",
     "Deep relationship across RPS, Bank Sweep, and O2 programs. Familiarity with Schwab's security, "
     "compliance, and governance standards. Faster ramp-up, lower delivery risk."),
    ("2. Deep SME Knowledge of Legacy Estates",
     "Strong understanding of COBOL mainframes, .NET Framework monoliths, and batch systems. "
     "Ability to translate legacy behavior into BRDs, DDD artifacts, and modern APIs."),
    ("3. Proven AI-Driven Modernization Capability",
     "Demonstrated use of VS Code + Copilot, Google ADK, LangGraph, and Custom Agents across the full SDLC. "
     "Vertex AI Search for grounded analysis. Spec-Driven Development for traceability."),
    ("4. Factory-Style, Repeatable Approach",
     "Agentic SDLC designed for repeatability, auditability, and scale. "
     "MCP integrations with GitHub, Jira/Confluence, CI/CD, PCF, and observability tools. "
     "Not a one-off POC—a model that scales across the entire portfolio."),
    ("5. Vertex AI Search & Spec-Driven Development Expertise",
     "Pioneered the integration of Vertex AI Search for domain knowledge grounding and "
     "Spec-Driven Development as the backbone of agent orchestration at Schwab."),
]
for title, desc in reasons:
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(1), Inches(1))
    # Just use the main textbox
    pass

# Re-do as single text box
# Remove the dummy textboxes
txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
first = True
for title, desc in reasons:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph()
    first = False
    p.space_before = Pt(10)
    _add_run(p, title, size=16, bold=True, color=DARK_BLUE)
    p2 = tf2.add_paragraph()
    _add_run(p2, desc, size=14, color=DARK_GRAY)

_add_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 16 – THANK YOU / Q&A
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
for ph in slide.placeholders:
    if ph.has_text_frame:
        ph.text_frame.clear()

# Centered thank you
shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
shp.fill.solid()
shp.fill.fore_color.rgb = DARK_BLUE
shp.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_add_run(p, "Thank You", size=44, bold=True, color=WHITE)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)
_add_run(p2, "Agentic SDLC at Schwab", size=28, bold=False, color=LIGHT_BLUE)
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)
_add_run(p3, "Questions & Discussion", size=22, bold=False, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\n✅ Deck saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
