"""
Bank Sweep Modernization – Pure Greenfield via Chat UI / Studio
All personas (BA, Architect, Developer, QA, DevOps) interact through
prompts in a Chat UI / Studio. No legacy code involved.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──────────────────────────────────────────────────────────────────────
TEMPLATE = r"C:\Users\banu.parasuraman\Downloads\mainframe-modernization\SchwabPlan.com API Modernization-V3.pptx"
OUTPUT   = r"C:\Users\banu.parasuraman\Downloads\mainframe-modernization\Bank_Sweep_Greenfield_Approach.pptx"

# ── Brand colours ──────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5C)
MID_BLUE   = RGBColor(0x00, 0x6B, 0xA6)
LIGHT_BLUE = RGBColor(0x00, 0x96, 0xD6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x00, 0x7A, 0x33)
ORANGE     = RGBColor(0xE8, 0x6C, 0x00)
TEAL       = RGBColor(0x00, 0x80, 0x80)
PURPLE     = RGBColor(0x6B, 0x21, 0xA8)
CHAT_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # Dark chat-style background
CHAT_USER  = RGBColor(0x00, 0x7A, 0xCC)   # User bubble colour
CHAT_AI    = RGBColor(0x2D, 0x2D, 0x44)   # AI response bubble

# ── Helpers ────────────────────────────────────────────────────────────────────
def _font(run, name="Segoe UI", size=18, bold=False, color=DARK_GRAY):
    run.font.name = name; run.font.size = Pt(size); run.font.bold = bold
    if color: run.font.color.rgb = color

def _add_run(para, text, name="Segoe UI", size=18, bold=False, color=DARK_GRAY):
    run = para.add_run(); run.text = text; _font(run, name, size, bold, color); return run

def _set_cell(cell, text, size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    _add_run(p, text, size=size, bold=bold, color=color)
    cell.text_frame.word_wrap = True

def _header_cell(cell, text, size=13):
    _set_cell(cell, text, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE

def _alt_row(cell, ri):
    if ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY

def _add_title_bar(slide, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.95))
    shp.fill.solid(); shp.fill.fore_color.rgb = DARK_BLUE; shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _add_run(p, "  " + text, size=28, bold=True, color=WHITE)

def _add_subtitle(slide, text, top=Inches(1.1)):
    txBox = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.5))
    _add_run(txBox.text_frame.paragraphs[0], text, size=20, bold=True, color=MID_BLUE)

def _add_body_text(slide, text, top=Inches(1.7), left=Inches(0.6),
                   width=Inches(12), height=Inches(5.2), size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_run(p, line, size=size, color=DARK_GRAY); p.space_after = Pt(6)
    return txBox

def _add_bullet_list(slide, items, top=Inches(1.7), left=Inches(0.6),
                     width=Inches(11.5), height=Inches(5.2), size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(item, tuple):
            _add_run(p, "●  " + item[0], size=size, bold=True, color=DARK_BLUE)
            _add_run(p, " – " + item[1], size=size, color=DARK_GRAY)
        else:
            _add_run(p, "●  " + item, size=size, color=DARK_GRAY)
    return txBox

def _add_table(slide, headers, rows, left=Inches(0.5), top=Inches(1.7),
               width=Inches(12.3), row_height=Inches(0.45)):
    cols = len(headers); num_rows = len(rows) + 1
    ts = slide.shapes.add_table(num_rows, cols, left, top, width, row_height * num_rows)
    tbl = ts.table; cw = int(width / cols)
    for ci in range(cols): tbl.columns[ci].width = cw
    for ci, h in enumerate(headers): _header_cell(tbl.cell(0, ci), h)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci); _set_cell(c, val, size=13, bold=(ci==0), color=DARK_GRAY); _alt_row(c, ri)
    return ts

def _box(slide, text, left, top, w, h, fill, tc=WHITE, sz=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _add_run(tf.paragraphs[0], text, size=sz, bold=True, color=tc)
    shp.text_frame.margin_left = Pt(6); shp.text_frame.margin_right = Pt(6)
    return shp

def _arrow_r(slide, l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = MID_BLUE; s.line.fill.background()

def _arrow_d(slide, l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = MID_BLUE; s.line.fill.background()

def _footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.5), Inches(0.35))
    _add_run(txBox.text_frame.paragraphs[0],
             "© Mphasis 2026  |  Proprietary and confidential  |  March 2026", size=10, color=MED_GRAY)

def _persona_badge(slide, text, left, top, color):
    return _box(slide, text, left, top, Inches(1.5), Inches(0.32), color, WHITE, 10)

def _chat_bubble(slide, persona, prompt, left, top, width):
    """Simulate a chat-style prompt bubble."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.9))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_run(p, f"💬 {persona}:  ", size=13, bold=True, color=MID_BLUE)
    _add_run(p, f'"{prompt}"', size=13, bold=False, color=DARK_GRAY)
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation(TEMPLATE)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId; prs.part.drop_rel(rId); del prs.slides._sldIdLst[0]
blank = prs.slide_layouts[2]; cover = prs.slide_layouts[0]

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 – COVER
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(cover)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
for ph in slide.placeholders:
    if "Title" in ph.name:
        _add_run(ph.text_frame.paragraphs[0], "Bank Sweep Modernization", size=36, bold=True, color=WHITE)
    elif "Subtitle" in ph.name:
        _add_run(ph.text_frame.paragraphs[0], "Pure Greenfield  •  All Personas via Chat UI / Studio", size=22, color=WHITE)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 – AGENDA
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Agenda")
_add_bullet_list(slide, [
    "Approach Overview – Pure Greenfield via Chat UI / Studio",
    "The Chat UI / Studio Concept",
    "End-to-End Flow – All Personas, All Prompts",
    "BA → Epics & Stories via Prompts",
    "Architect → Design & API Specs via Prompts",
    "Developer → Code Generation via Prompts",
    "QA → Test Suite Generation via Prompts",
    "DevOps → CI/CD & Deployment via Prompts",
    "Prompt-Driven SDLC – Persona Summary",
    "Advantages of the Chat UI / Studio Approach",
], top=Inches(1.3), size=18)
_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 – APPROACH OVERVIEW
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Approach Overview – Pure Greenfield via Chat UI / Studio")

body = (
    "Bank Sweep is built entirely from scratch – no legacy code, no conversion, no migration.\n\n"
    "Every persona on the team works through a Chat UI / Studio interface:\n"
    "prompts in → artifacts out.\n\n"
    "    💬  BA prompts → Epics, stories, and acceptance criteria in Jira\n"
    "    💬  Architect prompts → Domain models, OpenAPI specs, ADRs\n"
    "    💬  Developer prompts → Spring Boot code, unit tests, PRs\n"
    "    💬  QA prompts → Test suites, BDD scenarios, test data\n"
    "    💬  DevOps prompts → CI/CD pipelines, deployment manifests, observability\n\n"
    "The Chat UI / Studio is the single interface for the entire SDLC.\n"
    "Each persona brings domain expertise; the AI agent does the heavy lifting."
)
_add_body_text(slide, body, top=Inches(1.15), size=16)

_box(slide, "No Legacy Code  •  No Conversion  •  Pure Greenfield  •  All Prompts via Chat UI / Studio",
     Inches(0.4), Inches(6.2), Inches(12.4), Inches(0.55), PURPLE, WHITE, 15)
_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 – THE CHAT UI / STUDIO CONCEPT
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "The Chat UI / Studio – Single Interface for the Entire SDLC")

# Central chat box
_box(slide, "Chat UI / Studio", Inches(4.5), Inches(1.3), Inches(4.3), Inches(0.7), CHAT_BG, WHITE, 20)

# Persona boxes around it
personas = [
    ("👤 Business\nAnalyst", Inches(0.3),  Inches(2.5), MID_BLUE),
    ("👤 Architect",         Inches(3.0),  Inches(2.5), TEAL),
    ("👤 Developer",         Inches(5.7),  Inches(2.5), GREEN),
    ("👤 QA Engineer",       Inches(8.4),  Inches(2.5), PURPLE),
    ("👤 DevOps",            Inches(11.1), Inches(2.5), ORANGE),
]
for label, l, t, c in personas:
    _box(slide, label, l, t, Inches(2.2), Inches(0.75), c, WHITE, 13)

# "All through prompts" arrow concept
_box(slide, "Each persona types natural-language prompts  →  AI Agent generates artifacts  →  Human reviews & approves",
     Inches(0.3), Inches(3.6), Inches(12.7), Inches(0.55), LIGHT_BLUE, DARK_BLUE, 14)

# Output artifacts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "What Comes Out of the Chat UI / Studio", size=18, bold=True, color=DARK_BLUE)

outputs = [
    ("BA", "Jira Epics, User Stories, Acceptance Criteria, BRDs"),
    ("Architect", "DDD Models, OpenAPI Specs, ADRs, Database Schemas, Sequence Diagrams"),
    ("Developer", "Spring Boot Code, Unit Tests, Integration Code, Pull Requests"),
    ("QA", "JUnit Suites, Gherkin/BDD Scenarios, Test Data, Coverage Reports"),
    ("DevOps", "Dockerfiles, K8s/PCF Manifests, CI/CD Configs, Monitoring Dashboards"),
]
for role, arts in outputs:
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    _add_run(p2, f"💬 {role}:  ", size=14, bold=True, color=MID_BLUE)
    _add_run(p2, arts, size=14, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 – END-TO-END FLOW
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "End-to-End Greenfield Flow – All Personas via Chat UI")

# Row 1
steps1 = [
    ("Step 1\nBA → Epics\n& Stories", MID_BLUE),
    ("Step 2\nArchitect →\nDesign & Specs", TEAL),
    ("Step 3\nDeveloper →\nCode Gen", GREEN),
]
bw = Inches(3.5); bh = Inches(1.1); sl = Inches(0.4); t1 = Inches(1.3)
for i, (lbl, clr) in enumerate(steps1):
    left = sl + i * Inches(4.15)
    _box(slide, lbl, left, t1, bw, bh, clr, WHITE, 14)
    if i < len(steps1)-1:
        _arrow_r(slide, left+bw+Inches(0.1), t1+Inches(0.35), Inches(0.5), Inches(0.35))

# Row 2
steps2 = [
    ("Step 4\nQA → Test\nSuite Gen", PURPLE),
    ("Step 5\nCode Review\n& PRs", MID_BLUE),
    ("Step 6\nDevOps → CI/CD\n& Deploy", DARK_BLUE),
]
t2 = Inches(2.9)
for i, (lbl, clr) in enumerate(steps2):
    left = sl + i * Inches(4.15)
    _box(slide, lbl, left, t2, bw, bh, clr, WHITE, 14)
    if i < len(steps2)-1:
        _arrow_r(slide, left+bw+Inches(0.1), t2+Inches(0.35), Inches(0.5), Inches(0.35))

_arrow_d(slide, Inches(12.3), t1+bh+Inches(0.05), Inches(0.3), Inches(0.35))

# Chat UI bar
_box(slide, "💬  Chat UI / Studio  –  The Single Interface for Every Step Above",
     Inches(0.4), Inches(4.4), Inches(12.4), Inches(0.55), CHAT_BG, WHITE, 15)

# How it works
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "How It Works at Every Step:", size=16, bold=True, color=DARK_BLUE)
flow_desc = [
    "1.  Persona opens the Chat UI / Studio",
    "2.  Types a natural-language prompt describing what they need",
    "3.  AI Agent generates the artifact (stories, code, tests, manifests, etc.)",
    "4.  Persona reviews, refines via follow-up prompts, and approves",
    "5.  Approved artifact is pushed to the target system (Jira, GitHub, CI/CD, etc.)",
]
for fd in flow_desc:
    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    _add_run(p2, fd, size=14, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 – BA: EPICS & STORIES
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 1 – BA → Epics & Stories via Prompts")
_persona_badge(slide, "👤 Business Analyst", Inches(11.0), Inches(0.3), MID_BLUE)

# Prompt → Output flow
_box(slide, "BA types prompt\nin Chat UI", Inches(0.4), Inches(1.4), Inches(3.0), Inches(0.8), MID_BLUE, WHITE, 13)
_arrow_r(slide, Inches(3.5), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "AI Agent generates\nEpics & Stories", Inches(4.1), Inches(1.4), Inches(3.0), Inches(0.8), TEAL, WHITE, 13)
_arrow_r(slide, Inches(7.2), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "BA reviews,\nrefines, approves", Inches(7.8), Inches(1.4), Inches(3.0), Inches(0.8), PURPLE, WHITE, 13)
_arrow_r(slide, Inches(10.9), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "→ Jira", Inches(11.5), Inches(1.4), Inches(1.4), Inches(0.8), DARK_BLUE, WHITE, 13)

# Example prompts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(12.3), Inches(4.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Example Prompts the BA Would Type", size=17, bold=True, color=DARK_BLUE)

prompts = [
    ("Define the epics", "\"Create epics for the Bank Sweep cash management module. Include sweep scheduling, threshold management, and account linking.\""),
    ("Break into stories", "\"Break the sweep scheduling epic into user stories with acceptance criteria. Use Given-When-Then format.\""),
    ("Add edge cases", "\"Add stories for error handling: what happens when a sweep fails mid-transfer? What about insufficient funds?\""),
    ("Refine acceptance", "\"Update the threshold management stories – thresholds should be configurable per account type, not global.\""),
    ("Push to Jira", "\"Create these epics and stories in Jira project BSWP, assign to Sprint 3 backlog.\""),
]
for label, prompt in prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(10)
    _add_run(p2, f"💬  {label}:", size=14, bold=True, color=MID_BLUE)
    p3 = tf.add_paragraph()
    _add_run(p3, f"    {prompt}", size=13, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 – ARCHITECT: DESIGN & SPECS
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 2 – Architect → Design & API Specs via Prompts")
_persona_badge(slide, "👤 Architect", Inches(11.4), Inches(0.3), TEAL)

# Flow
_box(slide, "Architect types\nprompt in Chat UI", Inches(0.4), Inches(1.4), Inches(3.0), Inches(0.8), TEAL, WHITE, 13)
_arrow_r(slide, Inches(3.5), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "AI Agent generates\nDesign Artifacts", Inches(4.1), Inches(1.4), Inches(3.0), Inches(0.8), MID_BLUE, WHITE, 13)
_arrow_r(slide, Inches(7.2), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "Architect reviews,\niterates, approves", Inches(7.8), Inches(1.4), Inches(3.0), Inches(0.8), PURPLE, WHITE, 13)
_arrow_r(slide, Inches(10.9), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "→ Repo", Inches(11.5), Inches(1.4), Inches(1.4), Inches(0.8), DARK_BLUE, WHITE, 13)

# Example prompts & outputs
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(5.8), Inches(4.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Example Prompts", size=17, bold=True, color=DARK_BLUE)

arch_prompts = [
    "\"Design a microservice architecture for Bank Sweep with bounded contexts for Scheduling, Accounts, and Transfers.\"",
    "\"Generate an OpenAPI 3.0 spec for the Sweep Scheduling service with CRUD operations and webhook notifications.\"",
    "\"Create a database schema for the Accounts service using PostgreSQL. Include audit columns and soft deletes.\"",
    "\"Write an ADR for choosing event-driven communication between Scheduling and Transfers services.\"",
    "\"Generate a sequence diagram for the end-to-end sweep execution flow.\"",
]
for prompt in arch_prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "💬  " + prompt, size=13, color=DARK_GRAY)

# Outputs
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.0))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_add_run(p, "Generated Artifacts", size=17, bold=True, color=DARK_BLUE)

arch_outputs = [
    ("DDD Bounded Contexts", "Aggregates, entities, value objects per service"),
    ("OpenAPI 3.0 Specs", "Full API contracts with schemas, examples, error models"),
    ("Database Schemas", "DDL scripts, ER diagrams, migration files"),
    ("ADRs", "Architecture Decision Records with rationale & trade-offs"),
    ("Sequence Diagrams", "Mermaid / PlantUML for key workflows"),
    ("Integration Contracts", "AsyncAPI specs for event-driven messaging"),
]
for label, desc in arch_outputs:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "●  " + label, size=14, bold=True, color=TEAL)
    p3 = tf2.add_paragraph()
    _add_run(p3, "    " + desc, size=13, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 – DEVELOPER: CODE GENERATION
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 3 – Developer → Code Generation via Prompts")
_persona_badge(slide, "👤 Developer", Inches(11.2), Inches(0.3), GREEN)

# Flow
_box(slide, "Dev picks Jira\nticket, opens\nChat UI", Inches(0.4), Inches(1.4), Inches(2.8), Inches(0.9), GREEN, WHITE, 12)
_arrow_r(slide, Inches(3.3), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "Prompts for\nimplementation", Inches(3.8), Inches(1.4), Inches(2.8), Inches(0.9), MID_BLUE, WHITE, 12)
_arrow_r(slide, Inches(6.7), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "AI generates\nSpring Boot code\n+ unit tests", Inches(7.2), Inches(1.4), Inches(2.8), Inches(0.9), TEAL, WHITE, 12)
_arrow_r(slide, Inches(10.1), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "Dev reviews\n→ creates PR", Inches(10.6), Inches(1.4), Inches(2.4), Inches(0.9), PURPLE, WHITE, 12)

# Example prompts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(5.8), Inches(4.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Example Prompts", size=17, bold=True, color=DARK_BLUE)

dev_prompts = [
    "\"Implement the SweepScheduler service per the OpenAPI spec. Use Spring Boot 3, Java 17, Clean Architecture.\"",
    "\"Generate the JPA entities for the Accounts schema. Add Lombok, auditing, and soft delete support.\"",
    "\"Write a REST controller for the Threshold Management API with validation and error handling.\"",
    "\"Create unit tests for the SweepExecutionService covering success, insufficient funds, and partial failure.\"",
    "\"Refactor the Transfers service to use async messaging via Spring Cloud Stream.\"",
]
for prompt in dev_prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "💬  " + prompt, size=13, color=DARK_GRAY)

# Outputs
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.7), Inches(5.8), Inches(4.0))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_add_run(p, "Generated Code Artifacts", size=17, bold=True, color=DARK_BLUE)

dev_outputs = [
    ("Spring Boot Services", "Controllers, services, repositories – Clean Architecture"),
    ("JPA Entities & DTOs", "Database entities, request/response DTOs, mappers"),
    ("Unit Tests", "JUnit 5 + Mockito, aligned to story acceptance criteria"),
    ("API Validation", "Bean validation, custom validators, error handlers"),
    ("Configuration", "application.yml, profiles, feature flags"),
    ("Pull Request", "Code pushed to branch, PR created with description"),
]
for label, desc in dev_outputs:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "●  " + label, size=14, bold=True, color=GREEN)
    p3 = tf2.add_paragraph()
    _add_run(p3, "    " + desc, size=13, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 – QA: TEST SUITE GENERATION
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 4 – QA → Test Suite Generation via Prompts")
_persona_badge(slide, "👤 QA Engineer", Inches(11.0), Inches(0.3), PURPLE)

# Flow
_box(slide, "QA reviews\nJira stories &\nOpenAPI specs", Inches(0.4), Inches(1.4), Inches(2.8), Inches(0.9), PURPLE, WHITE, 12)
_arrow_r(slide, Inches(3.3), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "Prompts for\ntest generation", Inches(3.8), Inches(1.4), Inches(2.8), Inches(0.9), MID_BLUE, WHITE, 12)
_arrow_r(slide, Inches(6.7), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "AI generates\ntest suites &\ntest data", Inches(7.2), Inches(1.4), Inches(2.8), Inches(0.9), TEAL, WHITE, 12)
_arrow_r(slide, Inches(10.1), Inches(1.65), Inches(0.4), Inches(0.3))
_box(slide, "QA reviews\n& executes", Inches(10.6), Inches(1.4), Inches(2.4), Inches(0.9), GREEN, WHITE, 12)

# Example prompts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(5.8), Inches(4.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Example Prompts", size=17, bold=True, color=DARK_BLUE)

qa_prompts = [
    "\"Generate Gherkin BDD scenarios for the sweep scheduling stories, covering happy path, edge cases, and failures.\"",
    "\"Create integration tests for the end-to-end sweep flow using TestContainers and Spring Boot Test.\"",
    "\"Generate test data sets for the Accounts service – include valid, boundary, and invalid data.\"",
    "\"Write API contract tests to validate the Scheduling service against its OpenAPI spec.\"",
    "\"Generate a test coverage report and identify gaps against the Jira acceptance criteria.\"",
]
for prompt in qa_prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "💬  " + prompt, size=13, color=DARK_GRAY)

# Outputs
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.7), Inches(5.8), Inches(4.0))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_add_run(p, "Generated Test Artifacts", size=17, bold=True, color=DARK_BLUE)

qa_outputs = [
    ("Gherkin / BDD Scenarios", "Given-When-Then for every acceptance criterion"),
    ("JUnit Integration Tests", "Cross-service workflow tests with TestContainers"),
    ("API Contract Tests", "OpenAPI spec validation (RestAssured / Pact)"),
    ("Test Data Sets", "Valid, boundary, negative – ready for automated runs"),
    ("Coverage Reports", "Story-level coverage mapping with gap analysis"),
    ("Performance Test Scripts", "JMeter / Gatling scripts for load testing"),
]
for label, desc in qa_outputs:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "●  " + label, size=14, bold=True, color=PURPLE)
    p3 = tf2.add_paragraph()
    _add_run(p3, "    " + desc, size=13, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 – CODE REVIEW via Chat UI
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 5 – Code Review & PRs via Chat UI")

_add_subtitle(slide, "All Reviewers Use Prompts to Analyze, Validate, and Approve", top=Inches(1.1))

# Left – Reviewer prompts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.8))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Reviewer Prompts in Chat UI", size=17, bold=True, color=DARK_BLUE)

review_prompts = [
    "\"Review this PR for code quality, security vulnerabilities, and SOLID violations.\"",
    "\"Check if this implementation matches the acceptance criteria in BSWP-142.\"",
    "\"Verify the Scheduling controller follows our OpenAPI spec – flag any deviations.\"",
    "\"Analyze test coverage for this PR – are all edge cases from the Gherkin scenarios covered?\"",
    "\"Generate a PR summary and changelog entry for this merge.\"",
    "\"Are there any performance concerns with the sweep batch query in this PR?\"",
]
for prompt in review_prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "💬  " + prompt, size=13, color=DARK_GRAY)

# Right – What the agent provides
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_add_run(p, "Agent Provides", size=17, bold=True, color=DARK_BLUE)

review_outputs = [
    ("Code Quality Report", "Violations, anti-patterns, security issues flagged"),
    ("Spec Compliance Check", "Line-by-line comparison to OpenAPI contract"),
    ("Test Coverage Analysis", "Maps tests to acceptance criteria, highlights gaps"),
    ("Suggested Improvements", "Refactoring suggestions with code examples"),
    ("Auto-Generated PR Summary", "What changed, why, and what to look for"),
    ("Performance Analysis", "Query plan analysis, N+1 detection, memory concerns"),
]
for label, desc in review_outputs:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "●  " + label, size=14, bold=True, color=MID_BLUE)
    p3 = tf2.add_paragraph()
    _add_run(p3, "    " + desc, size=13, color=DARK_GRAY)

# Bottom note
_box(slide, "Human reviewer makes the final approve / request-changes decision. The Chat UI is the reviewer's co-pilot.",
     Inches(0.4), Inches(6.3), Inches(12.4), Inches(0.5), MID_BLUE, WHITE, 14)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 – DEVOPS: CI/CD & DEPLOYMENT
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Step 6 – DevOps → CI/CD & Deployment via Prompts")
_persona_badge(slide, "👤 DevOps", Inches(11.4), Inches(0.3), ORANGE)

# Flow
_box(slide, "DevOps defines\ndeployment needs\nin Chat UI", Inches(0.4), Inches(1.4), Inches(3.0), Inches(0.8), ORANGE, WHITE, 12)
_arrow_r(slide, Inches(3.5), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "AI generates\npipeline & infra\nconfigs", Inches(4.1), Inches(1.4), Inches(3.0), Inches(0.8), MID_BLUE, WHITE, 12)
_arrow_r(slide, Inches(7.2), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "DevOps reviews,\ntests in staging", Inches(7.8), Inches(1.4), Inches(3.0), Inches(0.8), TEAL, WHITE, 12)
_arrow_r(slide, Inches(10.9), Inches(1.6), Inches(0.5), Inches(0.3))
_box(slide, "→ Prod\n(PCF/K8s)", Inches(11.5), Inches(1.4), Inches(1.4), Inches(0.8), DARK_BLUE, WHITE, 12)

# Prompts
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(5.8), Inches(4.0))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
_add_run(p, "Example Prompts", size=17, bold=True, color=DARK_BLUE)

devops_prompts = [
    "\"Generate a multi-stage Dockerfile for the Scheduling service. Optimize for minimal image size.\"",
    "\"Create a GitHub Actions CI/CD pipeline: build, test, scan, deploy to staging, then prod.\"",
    "\"Generate Kubernetes deployment manifests with HPA, resource limits, and readiness probes.\"",
    "\"Set up Micrometer metrics and Spring Boot Actuator health checks for all services.\"",
    "\"Create a rollback runbook for the Bank Sweep deployment on PCF.\"",
    "\"Generate Grafana dashboard JSON for monitoring sweep execution latency and error rates.\"",
]
for prompt in devops_prompts:
    p2 = tf.add_paragraph(); p2.space_before = Pt(7)
    _add_run(p2, "💬  " + prompt, size=13, color=DARK_GRAY)

# Outputs
txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.6), Inches(5.8), Inches(4.0))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p = tf2.paragraphs[0]
_add_run(p, "Generated DevOps Artifacts", size=17, bold=True, color=DARK_BLUE)

devops_outputs = [
    ("Dockerfiles", "Multi-stage, optimized, with security scanning"),
    ("CI/CD Pipelines", "GitHub Actions / Jenkins – build, test, deploy"),
    ("K8s / PCF Manifests", "Deployments, services, HPA, config maps"),
    ("Observability", "Actuator, Micrometer, Grafana dashboards"),
    ("Runbooks", "Deployment, rollback, incident response procedures"),
    ("Infrastructure as Code", "Terraform / Helm charts for environment setup"),
]
for label, desc in devops_outputs:
    p2 = tf2.add_paragraph(); p2.space_before = Pt(8)
    _add_run(p2, "●  " + label, size=14, bold=True, color=ORANGE)
    p3 = tf2.add_paragraph()
    _add_run(p3, "    " + desc, size=13, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 – PERSONA SUMMARY TABLE
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Prompt-Driven SDLC – All Personas via Chat UI / Studio")

headers = ["Persona", "SDLC Phase", "What They Prompt For", "Output Artifacts", "Target System"]
rows = [
    ["Business\nAnalyst",
     "Requirements",
     "Epics, stories, acceptance\ncriteria, edge cases",
     "Jira epics & stories\nwith Given-When-Then",
     "Jira"],
    ["Architect",
     "Design",
     "Domain models, API specs,\nDB schemas, ADRs",
     "OpenAPI specs, DDL,\nDDD models, diagrams",
     "GitHub Repo"],
    ["Developer",
     "Implementation",
     "Service code, entities,\ncontrollers, unit tests",
     "Spring Boot code,\nJUnit tests, PRs",
     "GitHub\n(branch + PR)"],
    ["QA Engineer",
     "Testing",
     "BDD scenarios, integration\ntests, test data, coverage",
     "Gherkin, JUnit suites,\ncoverage reports",
     "GitHub +\nCI pipeline"],
    ["Reviewer\n(Dev/Lead)",
     "Code Review",
     "Quality analysis, spec\ncompliance, coverage gaps",
     "Review comments,\nPR summaries",
     "GitHub PR"],
    ["DevOps",
     "Deployment",
     "Dockerfiles, CI/CD,\nK8s manifests, monitoring",
     "Pipeline configs,\nobservability setup",
     "CI/CD +\nPCF / K8s"],
]
_add_table(slide, headers, rows, top=Inches(1.15), row_height=Inches(0.8))

_box(slide, "One interface (Chat UI / Studio)  •  Six personas  •  Full SDLC  •  All through natural-language prompts",
     Inches(0.4), Inches(6.2), Inches(12.4), Inches(0.5), PURPLE, WHITE, 14)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 – ADVANTAGES
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Advantages of the Chat UI / Studio Approach")

advantages = [
    ("Single Interface for Everyone",
     "BAs, architects, developers, QA, and DevOps all use the same Chat UI. No tool switching, "
     "no context loss. One interface, one conversation history, one source of truth."),
    ("Natural Language = Low Barrier",
     "No one needs to learn new tools or DSLs. Every persona works in their natural language. "
     "Domain experts can contribute directly without engineering translation."),
    ("Iterative Refinement via Follow-Up Prompts",
     "Don't like the output? Keep prompting. The chat UI supports multi-turn conversations, "
     "so every artifact is refined until the persona approves it."),
    ("Speed: Minutes, Not Days",
     "BRDs that took days → generated in minutes. OpenAPI specs that took a week → generated in one conversation. "
     "Code that took sprints → generated and reviewed in hours."),
    ("Full Traceability",
     "Every artifact traces back to the prompt that created it. The conversation history is a complete "
     "audit trail of why every design, code, and test decision was made."),
    ("Human Expertise + AI Execution",
     "Personas bring domain knowledge, judgment, and approval authority. "
     "AI brings speed, consistency, and pattern knowledge. The best of both."),
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.7))
tf = txBox.text_frame; tf.word_wrap = True
first = True
for title, desc in advantages:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(10)
    _add_run(p, "●  " + title, size=16, bold=True, color=DARK_BLUE)
    p2 = tf.add_paragraph()
    _add_run(p2, "    " + desc, size=14, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 14 – KEY TAKEAWAYS
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()
_add_title_bar(slide, "Key Takeaways")

takeaways = [
    ("Pure Greenfield – No Legacy Code",
     "Bank Sweep is built entirely from scratch. No COBOL, no conversion, no migration. "
     "Clean, modern Spring Boot from day one."),
    ("Chat UI / Studio Is the SDLC",
     "Every persona – BA, Architect, Developer, QA, DevOps – works through the same chat interface. "
     "Prompts in, artifacts out, review and approve."),
    ("Prompt → Artifact → Review → Ship",
     "The universal workflow: type a prompt, get a generated artifact, review it, refine via follow-ups, "
     "then push to the target system (Jira, GitHub, CI/CD)."),
    ("Human Judgment at Every Step",
     "AI generates, humans approve. No artifact ships without persona review. "
     "The chat history provides full traceability and audit trail."),
    ("Massive Acceleration",
     "What traditionally takes weeks (BRDs, design docs, code, test suites) is produced in "
     "conversational sessions. Quality maintained through iterative refinement."),
]

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))
tf = txBox.text_frame; tf.word_wrap = True
first = True
for title, desc in takeaways:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(12)
    _add_run(p, "●  " + title, size=16, bold=True, color=DARK_BLUE)
    p2 = tf.add_paragraph()
    _add_run(p2, "    " + desc, size=14, color=DARK_GRAY)

_footer(slide)

# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 15 – THANK YOU
# ──────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
for ph in slide.placeholders:
    if ph.has_text_frame: ph.text_frame.clear()

shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
shp.fill.solid(); shp.fill.fore_color.rgb = DARK_BLUE; shp.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(2), Inches(2.0), Inches(9), Inches(4))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_add_run(p, "Thank You", size=44, bold=True, color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)
_add_run(p2, "Bank Sweep Modernization", size=28, color=LIGHT_BLUE)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(10)
_add_run(p3, "Pure Greenfield  •  All Personas via Chat UI / Studio", size=20, color=WHITE)
p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(30)
_add_run(p4, "Questions & Discussion", size=22, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"\n✅ Deck saved to: {OUTPUT}")
print(f"   Total slides: {len(prs.slides)}")
